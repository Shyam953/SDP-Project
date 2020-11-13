from django.shortcuts import render
from django.conf import settings
from django.core.files.storage import FileSystemStorage
from django.http import HttpResponse
from django.views.decorators.csrf import csrf_exempt
from django.http import FileResponse
import json
import mimetypes
import os
from django.utils.encoding import smart_str 

#XML Parser 
import xml.etree.ElementTree as ET

#PPT creation
#for presentation
from pptx import Presentation
#Text Color
from pptx.dml.color import RGBColor
#Pt for making bullets, Inches for provide image width & hight in inches.
from pptx.util import Pt,Inches

#blob & Image
import base64
from PIL import Image
from io import BytesIO

# Create your views here.

def simple_upload(request):
    return render(request, 'simple_upload.html')

def about_us(request):
    return render(request,'about.html')


@csrf_exempt
def edit_lists(request):
    if request.is_ajax:
        tasks = request.POST['data_arr']
        dt = json.loads(tasks)
        # write content of data_arr in file a.txt
        f = open('media/a.txt', 'w',)
        f.write('\n'.join(dt))
        f.close()
       
        titlename=""
        #parsing XML content of a.txt file
        tree=ET.parse(os.path.join(settings.MEDIA_ROOT, 'a.txt'))
        root = tree.getroot()
        tag = root.tag

        prs = Presentation()
        bullet_slide_layout = prs.slide_layouts[1] # 1 - Title and Content
        # 0 - Title (presentation title slide)
        # 1 - Title and Content
        # 2 - Section Header (sometimes called Segue)
        # 3 - Two Content (side by side bullet textboxes)
        # 4 - Comparison (same but additional title for each side by side content box)
        # 5 - Title Only
        # 6 - Blank
        # 7 - Content with Caption
        # 8 - Picture with Caption

        for title in root:
            c = title.get('id')
            
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes #shpae contains Text/paragaphs/chart/table etc
            
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]   #containing sequence of placeholder shapes in this slide.
            
            title_shape.text = c
            text_frame = title_shape.text_frame # A shape may have a text frame that can contain display text for that we need textframe
            p = text_frame.paragraphs[0]    # p for paragraphs
            p.text = c
            font = p.font
            font.bold = True
            
            p.space_before = 0 # if any spaces before paragaphs it remove that spaces.
            k = 0
            for mp in title:
                img_list = mp.findall('img')
                img_count = len(img_list)
                marg = 0
                c = mp.get('id')
                tf = body_shape.text_frame
                if(k==0):
                    tf.text = c
                    k=1
                else:
                    p = tf.add_paragraph()
                    p.text += c
                    font = p.font
                    
                for sp in mp:
                    if(sp.tag=='sp'):
                        p = tf.add_paragraph()
                        p.level = 1 
                        p.text += str(sp.text)
                        font = p.font
                        font.size = Pt(20)
                        font.color.rgb = RGBColor(0x00, 0x00, 0x80)
                        
                #In one slide of PPT we contains maximum 2 images in one slide...
                
                for img in mp.findall('img'):
                    data = img.get('id')
                    data = data.split(",")  #data:image/png;base64 , iVBORw0K ..... we need only second part ie,iVBORw0K ....
                    data = data[1]
                    im = Image.open(BytesIO(base64.b64decode(data))) #iVBORw0.... read this string and make a approproate image  
                    im.save(os.path.join(settings.MEDIA_ROOT,'imgpdf2ppt.png'), 'PNG') #save('path','Format') format ie, jpg,png...
                    img_path = os.path.join(settings.MEDIA_ROOT,'imgpdf2ppt.png') 
                    if(img_count==1):
                        pic = slide.shapes.add_picture(img_path, Inches(3), Inches(4.3),width=Inches(4), height=Inches(3))
                    elif(img_count==2):
                        pic = slide.shapes.add_picture(img_path, Inches(marg+0.5), Inches(4.5),width=Inches(4), height=Inches(3))
                        marg = marg + 4.5

        # get first title for PPT name               
        for title in root: 
            titlename = (title.get('id'))
            print(titlename)
            break
            
        # save ppt at particular location
        filename = titlename.strip() + '.pptx'
        prs.save(os.path.join(settings.MEDIA_ROOT,filename))
        #return file in response for download
        return HttpResponse(filename)

        
